# -*- coding: utf-8 -*-
import os
import uuid
import re
import logging
from flask import Flask, request, render_template, send_file, jsonify, after_this_request
from werkzeug.utils import secure_filename
import requests
import json
from docx import Document as DocxDocument
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt
# Removed unused MSO_AUTO_SIZE
from pptx.enum.text import MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
from dotenv import load_dotenv

# Load environment variables from a .env file if it exists
load_dotenv()

# --- Configuration ---
UPLOAD_FOLDER = 'uploads'
GENERATED_FOLDER = 'generated'
ALLOWED_EXTENSIONS = {'docx', 'pdf'}
TEMPLATES = {
    'professional': {
        'title_color': RGBColor(0, 32, 96), 'accent_color': RGBColor(0, 112, 192), 'background_color': RGBColor(255, 255, 255), 'text_color': RGBColor(51, 51, 51), 'notes_suggestion_color': RGBColor(0, 100, 0)
    },
    'creative': {
        'title_color': RGBColor(113, 43, 124), 'accent_color': RGBColor(216, 101, 79), 'background_color': RGBColor(249, 245, 237), 'text_color': RGBColor(64, 64, 64), 'notes_suggestion_color': RGBColor(113, 43, 124)
    },
    'minimalist': {
        'title_color': RGBColor(64, 64, 64), 'accent_color': RGBColor(128, 128, 128), 'background_color': RGBColor(255, 255, 255), 'text_color': RGBColor(51, 51, 51), 'notes_suggestion_color': RGBColor(80, 80, 80)
    }
}
AZURE_ENDPOINT = os.environ.get("AZURE_OPENAI_ENDPOINT")
AZURE_API_KEY = os.environ.get("AZURE_OPENAI_API_KEY")

# --- Flask App Setup ---
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['GENERATED_FOLDER'] = GENERATED_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 32 * 1024 * 1024 # 32 MB max upload size
app.config['DEFAULT_TEMPLATE'] = 'professional'
os.makedirs(UPLOAD_FOLDER, exist_ok=True); os.makedirs(GENERATED_FOLDER, exist_ok=True)
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# --- Helper Functions ---
def allowed_file(filename):
    """Checks if the filename has an allowed extension."""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_docx(filepath):
    """Extracts text from a DOCX file."""
    try:
        doc = DocxDocument(filepath)
        full_text = [p.text for p in doc.paragraphs if p.text.strip()]
        return '\n'.join(full_text)
    except Exception as e:
        logging.error(f"Error extracting text from DOCX '{filepath}': {e}", exc_info=True)
        raise ValueError(f"Could not process DOCX file: {e}")

def extract_text_from_pdf(filepath):
    """Extracts text from a PDF file."""
    # NOTE: Requires Python 3.8+ for walrus operator :=
    try:
        reader = PdfReader(filepath)
        full_text = []
        for page in reader.pages:
            if (text := page.extract_text()) and text.strip():
                full_text.append(text.strip())
        return '\n'.join(full_text)
    except Exception as e:
        logging.error(f"Error extracting text from PDF '{filepath}': {e}", exc_info=True)
        raise ValueError(f"Could not process PDF file: {e}")

# *** build_llm_prompt FUNCTION ***
def build_llm_prompt(document_text, template_name='professional', audience="", tone=""):
    """Builds the detailed prompt for the LLM."""
    max_chars = 400000 # Limit input text size (adjust as needed for token limits)
    truncated_text = document_text[:max_chars]
    if len(document_text) > max_chars:
        logging.warning(f"Input document text truncated to {max_chars} characters for LLM prompt.")

    context_lines = []
    if audience: context_lines.append(f"**Target Audience:** {audience}")
    if tone: context_lines.append(f"**Desired Tone:** {tone}")
    if template_name: context_lines.append(f"**Visual Style Hint:** '{template_name.capitalize()}' theme.")
    context_str = "\n".join(context_lines)

    # The detailed prompt structure remains the same
    prompt = f"""
You are an expert presentation designer AND coach creating a **robust and detailed first draft** PowerPoint outline (~10-15 slides) with speaker notes and suggestions. Structure information logically for the audience/tone. Ensure bullet points are informative (~10-15 words/short sentences).

**User-Provided Context:**
{context_str if context_lines else "**User Context:** None provided."}

**CRITICAL INSTRUCTIONS:**
1. Generate the entire response following the structure below for **EVERY** slide.
2. **ALL fields listed (Slide Title, Content Type, Key Message, Bullets (at least one), Visual Suggestion, Design Note, Notes, Elaboration, Enhancement Suggestion, Best Practice Tip) ARE REQUIRED for each slide block.** Do not omit fields. Provide meaningful content or state 'None' or 'N/A' where appropriate but the field label MUST be present.
3. Start each slide block *exactly* with `---` on its own line.
4. Ensure bullets start with '- ' and contain substantial information, not just keywords.
5. Elaboration MUST expand significantly on the bullets/key message for the speaker.
6. Enhancement Suggestion and Best Practice Tip MUST be actionable and relevant.

**Required Slide Block Format (Example Included):**

---
Slide Title: [Concise Title in Title Case]
Content Type: [Text Only / Text and Image / Text and Chart / etc.]
Key Message: [**Single sentence essence** tailored for Audience/Tone.]
- [Bullet 1: **Informative point** (~10-15 words/short sentence) directly from or synthesized from the source text.]
- [Bullet 2: (Optional) ditto. Use 3-5 relevant bullets per slide.]
- [...]
Visual Suggestion: [Specific & actionable suggestion (e.g., 'Bar chart comparing Q1 vs Q2 sales', 'Icon representing collaboration') or "Text Focus".]
Design Note: [Optional: Suggest specific emphasis, layout idea, or callout.]
Notes: [Optional: Brief background context, source reference, or data point origin.]
**Elaboration:** [REQUIRED: Expand significantly on the slide's points (2-4 detailed sentences). Provide context, nuance, or talking points for the speaker.]
**Enhancement Suggestion:** [REQUIRED: Offer 1-2 concrete, actionable ideas for the user to improve this specific slide (e.g., 'Add a customer quote here', 'Break this into two slides if time permits').]
**Best Practice Tip:** [REQUIRED: Provide one relevant presentation best practice tip for this type of slide content (e.g., 'Use high-contrast colors for accessibility', 'Limit text to 6 lines per slide').]
---

**Example of ONE complete slide block:**
---
Slide Title: Understanding the Core Problem
Content Type: Text Only
Key Message: Current manual processes lead to significant delays and potential errors in reporting, impacting timely decisions.
- Manual data entry is time-consuming, taking approximately 4 hours per week per analyst involved in the process.
- Lack of real-time validation mechanisms increases the risk of inaccurate financial statements being circulated.
- Reporting delays directly impact the ability of executive leadership to make informed, timely strategic decisions.
Visual Suggestion: Icon representing 'Time Wasted' or 'Error Symbol'. Consider a simple process flow showing the manual steps.
Design Note: Use bold text or a distinct color for the time/delay figures mentioned in the elaboration.
Notes: Data based on internal Q3 process review document and interviews with the finance team.
**Elaboration:** The 4-hour figure represents an average across five analysts; peak times near month-end are higher. Errors identified later require significant rework, sometimes delaying the final month-end close by up to two business days. Leadership relies on these reports for the Wednesday morning strategy meeting, requiring accuracy by EOD Tuesday.
**Enhancement Suggestion:** Quantify the potential financial cost of errors (e.g., average cost per correction, impact of delayed decisions) if possible. Consider adding a brief, anonymous quote from an analyst about the manual process friction.
**Best Practice Tip:** When presenting a problem slide, clearly articulate the 'so what?' – the direct impact on key business objectives or stakeholders (like leadership decision-making).
---

**(Continue generating ALL subsequent slide blocks using the EXACT format above, including ALL required fields)**

**Presentation Flow Guidance (Adapt as needed based on content):**
1. Title Slide, 2. Agenda/Overview, 3. Introduction/Problem Statement, 4. Section 1 (1-3 slides), 5. Section 2 (1-3 slides), 6. Section 3 (1-3 slides) [Adjust section count based on text length/complexity], 7. Key Findings/Analysis (if applicable), 8. Recommendations/Solutions, 9. Next Steps/Call to Action, 10. Conclusion/Summary, 11. Q&A/Contact Information

**Content & Style Guidelines:**
* Structure Comprehensively: Ensure a logical flow from start to finish.
* Informative Bullets: ~10-15 words per bullet, focus on clarity and impact.
* Prioritize Key Info: Extract the most important messages from the source text.
* Tailor Tone/Audience: Reflect the specified audience and tone in language and focus.
* Elaborate Meaningfully: Notes should add real value for the speaker.
* Actionable Suggestions: Enhancements and Tips should be practical.
* Data Storytelling: If data is present, weave it into a narrative.
* Handle Missing Content Gracefully: If source text is sparse for a section, indicate this clearly (e.g., in Notes) but still attempt to generate a placeholder slide with suggestions.

**Source Document Text:**
\"\"\"
{truncated_text}
\"\"\"

Generate the **complete** presentation outline now, following ALL instructions meticulously for **EVERY** slide. Ensure all required fields (Slide Title, Content Type, Key Message, Bullets, Visual Suggestion, Design Note, Notes, Elaboration, Enhancement Suggestion, Best Practice Tip) are present in each block starting with `---`.
"""
    return prompt

# *** call_llm FUNCTION (WITH FIX INCORPORATED) ***
def call_llm(prompt):
    if not AZURE_ENDPOINT or not AZURE_API_KEY:
        raise ValueError("AI Service endpoint or API key not configured in environment variables.") # More specific error

    headers = {
        "Content-Type": "application/json",
        "api-key": AZURE_API_KEY
    }
    # Consider making these configurable or constants
    max_output_tokens = 8192
    timeout_seconds = 300

    payload = {
        "messages": [
            {
                "role": "system",
                "content": "You are an expert presentation designer/coach creating detailed PowerPoint outlines following strict formatting rules, including mandatory elaboration and suggestion fields."
            },
            {
                "role": "user",
                "content": prompt
            }
        ],
        "max_tokens": max_output_tokens,
        "temperature": 0.5, # Adjust temperature as needed (0.5 is balanced)
        "top_p": 0.95
        # Consider adding other parameters like 'frequency_penalty', 'presence_penalty' if needed
    }

    try:
        logging.info(f"Calling Azure OpenAI API at {AZURE_ENDPOINT} (max_tokens: {max_output_tokens}, timeout: {timeout_seconds}s)...")
        response = requests.post(AZURE_ENDPOINT, headers=headers, json=payload, timeout=timeout_seconds)
        response.raise_for_status() # Raises HTTPError for bad responses (4xx or 5xx)

        result = response.json() # Parse JSON response

        # --- Robust response validation ---
        if not isinstance(result, dict):
             raise ValueError(f"Invalid response type from AI service. Expected dict, got {type(result)}.")

        choices = result.get('choices')
        if not isinstance(choices, list) or not choices:
            logging.error("LLM response missing 'choices' list or list is empty: %s", result)
            raise ValueError("Invalid response structure from AI service: 'choices' are missing or empty.")

        first_choice = choices[0]
        if not isinstance(first_choice, dict):
            raise ValueError(f"Invalid 'choices' item type. Expected dict, got {type(first_choice)}.")

        message = first_choice.get('message')
        if not isinstance(message, dict):
            logging.error("LLM response 'choice' missing 'message' dict: %s", first_choice)
            raise ValueError("Invalid response structure: 'message' is missing in choice.")

        llm_output = message.get('content')
        if not isinstance(llm_output, str):
            logging.error("LLM response 'message' missing 'content' string: %s", message)
            raise ValueError("Invalid response structure: 'content' is missing or not a string.")
        # --- End validation ---

        llm_output = llm_output.strip() # Strip leading/trailing whitespace
        finish_reason = first_choice.get('finish_reason', 'unknown')

        if finish_reason == 'length':
            logging.warning(f"LLM response may have been truncated because it reached the maximum token limit ({max_output_tokens} tokens).")
        elif finish_reason != 'stop':
             logging.warning(f"LLM response finished with reason: '{finish_reason}'. Expected 'stop'.")


        # Basic sanity check on content structure
        if not llm_output or "---" not in llm_output or "Slide Title:" not in llm_output:
            logging.warning(f"LLM output seems to lack the expected basic slide structure (--- and Slide Title:). Preview: {llm_output[:200]}...")
            # Decide if this should be a hard error or just a warning based on requirements
            # raise ValueError("AI response content lacks expected structure.")

        logging.info(f"LLM Response Received Successfully (finish_reason: {finish_reason}).")
        return llm_output

    except requests.exceptions.Timeout:
        logging.error(f"LLM API call timed out after {timeout_seconds} seconds.")
        raise ValueError(f"Request timed out after {timeout_seconds}s waiting for AI service.") # More user-friendly timeout message

    except requests.exceptions.RequestException as e:
        # --- FIX: Corrected try-except block for error details ---
        status_code = e.response.status_code if e.response is not None else "N/A"
        error_details_str = str(e) # Default error details to the exception string

        if e.response is not None:
            # Try to get more specific error details from the response body
            try:
                # Attempt to parse response body as JSON for detailed error messages
                error_data = e.response.json()
                # Look for common error structures (adapt based on Azure API specifics)
                if isinstance(error_data, dict) and 'error' in error_data:
                     error_details_str = json.dumps(error_data['error'])
                else:
                     error_details_str = json.dumps(error_data) # Convert full JSON details to string
            except json.JSONDecodeError:
                # If response is not JSON, use the raw text (limit length)
                error_details_str = e.response.text[:500] + ('...' if len(e.response.text) > 500 else '')
            except Exception as json_ex:
                 # Catch unexpected errors during JSON parsing/dumping
                 logging.warning(f"Could not parse error response body, using raw text. Error: {json_ex}")
                 error_details_str = e.response.text[:500] + ('...' if len(e.response.text) > 500 else '') # Fallback to raw text

        logging.error(f"LLM API request failed: {e}. Status Code: {status_code}. Response Details: {error_details_str}", exc_info=False) # Log details, but maybe not full traceback unless debugging

        # Provide more specific user-facing errors based on status code
        if status_code == 401:
            raise ValueError("AI service authentication failed. Please check your API key configuration.")
        elif status_code == 404:
             raise ValueError("AI service endpoint not found. Please check the AZURE_OPENAI_ENDPOINT URL.")
        elif status_code == 429:
            raise ValueError("AI service rate limit exceeded or quota reached. Please wait and try again later or check your service plan.")
        elif status_code == 400:
            # Include details if available, otherwise a generic message
            detail_msg = f": {error_details_str}" if error_details_str != str(e) else ""
            raise ValueError(f"AI service rejected the request due to invalid input (Bad Request){detail_msg}. Please check the generated prompt or input document.")
        elif status_code >= 500:
             raise ValueError(f"The AI service encountered an internal server error (Status: {status_code}). Please try again later.")
        else:
            # Generic error for other client/server issues
            raise ValueError(f"API communication error occurred (Status Code: {status_code}). Check network connection and API status.")
        # --- End Fix ---

    except json.JSONDecodeError as e:
         # Handle cases where the successful response (2xx) is not valid JSON
         logging.error(f"Failed to decode valid JSON response from LLM API: {e}", exc_info=True)
         # Log response text if possible
         response_text_preview = "N/A"
         try:
             if 'response' in locals() and hasattr(response, 'text'):
                 response_text_preview = response.text[:500] + ('...' if len(response.text) > 500 else '')
         except Exception as E:
              logging.warning(f"Could not get response text preview after JSONDecodeError: {E}")
         logging.error(f"LLM Raw Response Text Preview: {response_text_preview}")
         raise ValueError("Received an invalid or malformed JSON response from the AI service.")

    except Exception as e:
        # Catch any other unexpected errors during the process
        logging.error(f"An unexpected error occurred during the LLM call: {e}", exc_info=True)
        raise ValueError(f"An unexpected error occurred while communicating with the AI service.")

# *** parse_llm_output FUNCTION ***
def parse_llm_output(llm_text):
    """Parses the structured text output from LLM into a list of slide dictionaries."""
    slides = []
    # Regex to split by '---' possibly surrounded by whitespace, ensuring it's on its own line
    slide_blocks = re.split(r'\n\s*---\s*\n', llm_text.strip())

    # Define all expected fields, including required and optional ones
    core_required_fields = {'title', 'content_type', 'key_message'}
    suggestion_fields = {'elaboration', 'enhancement_suggestion', 'best_practice_tip'}
    optional_fields = {'visual', 'design_note', 'notes', 'bullets'} # Bullets are technically optional if content dictates
    all_expected_fields = core_required_fields.union(suggestion_fields).union(optional_fields)

    # Default text for required suggestion fields if LLM omits them (though prompt asks it not to)
    default_suggestion = "Suggestion not provided by AI."

    # Prefixes for easier parsing (lowercase for case-insensitivity)
    # Ensure keys match the field names used in the dictionaries
    prefixes = {
        'slide title:': ('title', 12),
        'content type:': ('content_type', 13),
        'key message:': ('key_message', 12),
        'visual suggestion:': ('visual', 18),
        'design note:': ('design_note', 12),
        'notes:': ('notes', 6),
        'elaboration:': ('elaboration', 12),
        'enhancement suggestion:': ('enhancement_suggestion', 23),
        'best practice tip:': ('best_practice_tip', 18),
    }

    for block_idx, block in enumerate(slide_blocks):
        block = block.strip()
        if not block or block == '---': # Skip empty blocks or blocks just containing the separator
            continue

        current_slide = {field: '' for field in all_expected_fields} # Initialize with defaults
        current_slide['bullets'] = [] # Initialize bullets as list
        current_field_key = None # Tracks the field being processed for multi-line content

        lines = block.split('\n')
        field_buffer = {} # Temporary storage for field content during line processing

        for line_num, line in enumerate(lines):
            line_strip = line.strip()
            line_lower = line.lower()
            if not line_strip: continue # Skip empty lines

            # --- Check for Bullets ---
            if line.startswith('- '):
                bullet_content = line[2:].strip()
                if bullet_content: # Only add non-empty bullets
                    field_buffer.setdefault('bullets', []).append(bullet_content)
                current_field_key = 'bullets' # Set context
                continue # Move to next line

            # --- Check for Field Prefixes ---
            matched_prefix = False
            for prefix, (field_key, prefix_len) in prefixes.items():
                if line_lower.startswith(prefix):
                    # Found a field prefix, store its content
                    field_content = line[prefix_len:].strip()
                    field_buffer[field_key] = field_content
                    current_field_key = field_key # Update context to the new field
                    matched_prefix = True
                    break # Stop checking prefixes for this line

            # --- Handle Multi-line Content ---
            if not matched_prefix and current_field_key and current_field_key in field_buffer:
                # If no prefix matched, and we know the current field context, append the line
                # This handles multi-line descriptions for fields like Elaboration, Notes, etc.
                # Exclude appending to 'bullets' field here, as bullets are handled separately
                if current_field_key != 'bullets':
                     field_buffer[current_field_key] += f"\n{line_strip}"

        # --- Finalize Slide Data after processing all lines in the block ---
        current_slide.update(field_buffer) # Update slide dict with buffered content

        # Validate Core Fields: Ensure they exist and are not empty after stripping
        core_fields_present_and_filled = all(
            f in current_slide and current_slide[f].strip() for f in core_required_fields
        )

        if core_fields_present_and_filled:
            # Apply defaults for suggestion fields if they are missing or empty
            for f in suggestion_fields:
                if not current_slide.get(f, '').strip():
                    current_slide[f] = default_suggestion

            # Ensure bullet list exists even if empty
            current_slide.setdefault('bullets', [])

            slides.append(current_slide)
        else:
            # Log skipped blocks more informatively
            missing_or_empty_core = [f for f in core_required_fields if not current_slide.get(f, '').strip()]
            title_preview = current_slide.get('title', 'N/A')
            logging.warning(
                f"Skipping Block {block_idx+1}. Missing/Empty Core Fields: {missing_or_empty_core}. "
                f"Title='{title_preview}'. Block Preview: '{block[:150]}...'"
            )

    if not slides:
        # Critical if no slides could be parsed
        logging.error(f"Failed to parse ANY valid slide blocks from the LLM output. Output Preview: {llm_text[:500]}...")
        raise ValueError("AI response parsing failed. Could not find any slides with the required structure (Title, Content Type, Key Message). Please check the AI's output format.")

    logging.info(f"Successfully parsed {len(slides)} slides from LLM output.")
    return slides


# *** create_presentation FUNCTION ***
def create_presentation(slides_data, output_path, template_name='professional'):
    """Creates a PowerPoint presentation from the parsed slide data."""
    template = TEMPLATES.get(template_name, TEMPLATES['professional']) # Fallback to default template
    prs = Presentation()
    # Set slide size to 16:9 aspect ratio (Widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Title Slide ---
    if slides_data:
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[0]) # Layout 0 is typically Title Slide
            title = slide.shapes.title # Get the title placeholder shape

            # Find the subtitle placeholder (often index 1, but check)
            subtitle = next((s for s in slide.placeholders if s.placeholder_format.idx == 1 and s != title), None)
            if subtitle is None and len(slide.placeholders) > 1:
                subtitle = slide.placeholders[1] # Fallback: Assume index 1 if specific check fails

            # Set Title Text
            title_text = slides_data[0].get('title', 'Presentation') # Default title
            if title and title.has_text_frame:
                tf = title.text_frame; tf.clear()
                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                run = p.add_run(); run.text = title_text
                run.font.color.rgb = template['title_color']
                run.font.size = Pt(44); run.font.bold = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE # Center vertically
                p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER # Center horizontally
            elif title: # Fallback if title shape exists but no text frame
                title.text = title_text

            # Set Subtitle Text (Using Key Message)
            if subtitle and subtitle.has_text_frame:
                tf = subtitle.text_frame; tf.clear()
                subtitle_text = slides_data[0].get('key_message', '').strip()
                if subtitle_text:
                    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                    run = p.add_run(); run.text = subtitle_text
                    run.font.color.rgb = template['accent_color']
                    run.font.size = Pt(24)
                    tf.vertical_anchor = MSO_ANCHOR.TOP # Align subtitle towards top
                    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                else:
                    tf.text = "" # Ensure empty if no text provided
            elif subtitle: # Fallback if subtitle shape exists but no text frame
                 try: subtitle.text = slides_data[0].get('key_message', '').strip()
                 except AttributeError: logging.warning("Subtitle placeholder on title slide lacks text frame and direct text setting failed.")

            # Add notes for the title slide
            notes_tf = slide.notes_slide.notes_text_frame; notes_tf.clear()
            add_formatted_notes(notes_tf, slides_data[0], template)

        except Exception as e:
            logging.error(f"Critical Error creating Title Slide: {e}", exc_info=True)
            # Consider adding a generic title slide if the first one fails critically
            try:
                 slide = prs.slides.add_slide(prs.slide_layouts[0])
                 title = slide.shapes.title
                 if title: title.text = "Presentation (Title Generation Error)"
            except Exception as fallback_e:
                 logging.error(f"Failed to add fallback title slide: {fallback_e}")

    # --- Content Slides ---
    # Use Layout 1 (Title and Content) as default for content slides
    content_slide_layout = prs.slide_layouts[1]

    for idx, slide_data in enumerate(slides_data[1:]): # Iterate starting from the second item
        slide_num_for_logging = idx + 2 # User-facing slide number (starts from 2)
        slide_title_for_error = slide_data.get('title', f'Untitled Slide {slide_num_for_logging}')

        try:
            slide = prs.slides.add_slide(content_slide_layout)
            title_shape = slide.shapes.title
            # Find the main content placeholder (often idx 1, but find largest if not title)
            content_placeholder = next((s for s in slide.placeholders if s.placeholder_format.idx == 1 and s != title_shape), None)
            if content_placeholder is None: # Fallback: find largest placeholder that isn't the title
                 non_title_placeholders = [p for p in slide.placeholders if p != title_shape]
                 if non_title_placeholders:
                     content_placeholder = max(non_title_placeholders, key=lambda x: x.width * x.height)

            # Set Slide Title
            if title_shape and title_shape.has_text_frame:
                title_text = slide_data.get('title', f'Slide {slide_num_for_logging}')
                tf = title_shape.text_frame; tf.clear()
                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                run = p.add_run(); run.text = title_text
                run.font.color.rgb = template['title_color']
                run.font.size = Pt(32); run.font.bold = True
            elif title_shape:
                title_shape.text = slide_data.get('title', f'Slide {slide_num_for_logging}')

            # Populate Content Placeholder
            if content_placeholder and content_placeholder.has_text_frame:
                tf = content_placeholder.text_frame; tf.clear()
                tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP

                # Add Key Message (Formatted differently, optional)
                key_message = slide_data.get('key_message', '').strip()
                if key_message:
                    p_key = tf.add_paragraph()
                    run_key = p_key.add_run(); run_key.text = key_message
                    run_key.font.bold = True; run_key.font.italic = True
                    run_key.font.color.rgb = template['accent_color']
                    run_key.font.size = Pt(20)
                    p_key.space_after = Pt(10) # Add space after key message

                # Add Bullets
                bullets = slide_data.get('bullets', [])
                for bullet_text in bullets:
                    bullet_text = bullet_text.strip()
                    if not bullet_text: continue
                    p_bullet = tf.add_paragraph(); p_bullet.text = bullet_text; p_bullet.level = 0
                    if p_bullet.runs: # Safely access the first run
                        run_bullet = p_bullet.runs[0]
                        run_bullet.font.size = Pt(18); run_bullet.font.color.rgb = template['text_color']
                    else:
                        logging.warning(f"Paragraph created without run for bullet: '{bullet_text}' on slide {slide_num_for_logging} ('{slide_title_for_error}')")

                # Add Visual Placeholder if suggested
                visual_suggestion = slide_data.get('visual', '').strip()
                # Check against common "no visual" phrases
                if visual_suggestion and visual_suggestion.lower() not in ["none", "none needed", "text focus", "n/a", "text only", "no visual needed", ""]:
                    # Define default position/size (e.g., right side) - adjust as needed
                    placeholder_left = Inches(7.0); placeholder_top = Inches(1.8)
                    placeholder_width = Inches(5.5); placeholder_height = Inches(4.5)
                    add_visual_placeholder(slide, visual_suggestion, template, placeholder_left, placeholder_top, placeholder_width, placeholder_height)

            elif content_placeholder:
                logging.warning(f"Content placeholder found on slide {slide_num_for_logging} ('{slide_title_for_error}') but it lacks a text frame.")
            else:
                logging.warning(f"No suitable content placeholder found on slide {slide_num_for_logging} ('{slide_title_for_error}'). Content may be missing.")

            # Add formatted notes to the notes slide page
            notes_tf = slide.notes_slide.notes_text_frame; notes_tf.clear()
            add_formatted_notes(notes_tf, slide_data, template)

        except Exception as e:
             # Log the error with specific slide context
             logging.error(f"Error processing content slide {slide_num_for_logging} ('{slide_title_for_error}'): {e}", exc_info=True)
             # Attempt to add an error placeholder slide to the presentation
             try:
                 # Use a blank or title-only layout for error slide
                 error_slide = prs.slides.add_slide(prs.slide_layouts[5]) # Layout 5 is often Title Only or Blank
                 title_shape = error_slide.shapes.title
                 if title_shape: title_shape.text = f"Error Generating Slide {slide_num_for_logging}"

                 # Add a textbox explaining the error
                 left, top, width, height = Inches(1), Inches(1.5), Inches(11), Inches(5)
                 txBox = error_slide.shapes.add_textbox(left, top, width, height)
                 tf_err = txBox.text_frame; tf_err.word_wrap = True
                 tf_err.text = f"Failed to generate content for slide:\n'{slide_title_for_error}'\n\nError Details:\n{str(e)[:500]}{'...' if len(str(e)) > 500 else ''}"
                 # Basic formatting for error text
                 if tf_err.paragraphs:
                     for p in tf_err.paragraphs:
                         if p.runs: p.runs[0].font.size = Pt(14)

                 logging.info(f"Added error placeholder slide for failed slide {slide_num_for_logging}")

             except Exception as inner_e:
                 logging.critical(f"CRITICAL FAILURE: Could not add error placeholder slide for slide {slide_num_for_logging}. Inner Exception: {inner_e}")
             continue # Important: Skip to the next slide data item

    try:
        prs.save(output_path)
        logging.info(f"Presentation saved successfully: {output_path}")
    except Exception as e:
        logging.error(f"Failed to save the final presentation file '{output_path}': {e}", exc_info=True)
        # Raise an IOError to signal failure in saving, handled in the route
        raise IOError(f"Could not save the PowerPoint file: {e}")


# *** HELPER FUNCTION for Formatted Notes ***
def add_formatted_notes(notes_text_frame, slide_data, template):
    """Adds structured and formatted notes to the notes slide page."""
    notes_tf = notes_text_frame; notes_tf.clear()
    default_suggestion_text = "Suggestion not provided by AI."
    suggestion_color = template.get('notes_suggestion_color', RGBColor(80, 80, 80)) # Default gray if not in template
    missing_color = RGBColor(128, 128, 128) # Color for "(Suggestion not generated...)" text

    def add_section(title, content_key, is_suggestion=False):
        content = slide_data.get(content_key, '').strip()
        if not content and not is_suggestion: return # Skip empty non-suggestion sections

        is_missing_suggestion = is_suggestion and (not content or content == default_suggestion_text)
        final_content = content if not is_missing_suggestion else "(Suggestion not generated by AI)"

        # Add Section Title
        p_title = notes_tf.add_paragraph()
        run_title = p_title.add_run(); run_title.text = title
        run_title.font.bold = True; run_title.font.size = Pt(11)
        if is_suggestion: run_title.font.color.rgb = suggestion_color
        p_title.space_after = Pt(2)

        # Add Content Paragraph
        p_content = notes_tf.add_paragraph()
        run_content = p_content.add_run(); run_content.text = final_content
        run_content.font.size = Pt(10)
        if is_suggestion:
            run_content.font.italic = True
            # Use specific colors for provided vs. missing suggestions
            run_content.font.color.rgb = missing_color if is_missing_suggestion else suggestion_color
        p_content.level = 1 # Indent content slightly
        p_content.space_after = Pt(6) # Space after content before next section/separator

    # Define the order and titles for notes sections
    notes_sections = [
        {"title": "Speaker Notes:", "key": "notes", "is_suggestion": False},
        {"title": "Elaboration / Talking Points:", "key": "elaboration", "is_suggestion": False}, # Treat elaboration as core content
        {"title": "💡 Enhancement Suggestion:", "key": "enhancement_suggestion", "is_suggestion": True},
        {"title": "⭐ Best Practice Tip:", "key": "best_practice_tip", "is_suggestion": True},
    ]

    for section in notes_sections:
        add_section(section["title"], section["key"], section["is_suggestion"])

    # Optional: Add a final separator if needed, but space_after might be sufficient
    # p_sep = notes_tf.add_paragraph(); run_sep = p_sep.add_run(); run_sep.text = "---"
    # run_sep.font.size = Pt(8); run_sep.font.color.rgb = RGBColor(180, 180, 180)


# *** add_visual_placeholder FUNCTION ***
def add_visual_placeholder(slide, visual_desc, template, left, top, width, height):
    """Adds a styled placeholder shape with the visual suggestion text."""
    try:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)

        # Style the placeholder shape
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240) # Light gray background
        shape.line.color.rgb = template.get('accent_color', RGBColor(128, 128, 128)) # Border uses accent color
        shape.line.width = Pt(1.5)
        shape.shadow.inherit = False # Disable default shadow

        # Configure text frame properties
        tf = shape.text_frame
        tf.margin_bottom = Inches(0.1); tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1); tf.margin_top = Inches(0.1)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE # Center text vertically
        tf.word_wrap = True
        tf.clear()

        # Add "Suggested Visual:" title
        p_title = tf.add_paragraph()
        run_title = p_title.add_run(); run_title.text = "Suggested Visual:"
        run_title.font.bold = True
        run_title.font.color.rgb = template.get('accent_color', RGBColor(0, 0, 0)) # Use accent color for title
        run_title.font.size = Pt(14)
        p_title.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        p_title.space_after = Pt(6) # Space between title and description

        # Add the actual visual description
        p_desc = tf.add_paragraph()
        run_desc = p_desc.add_run()
        # Truncate long descriptions to fit reasonably
        run_desc.text = visual_desc[:250] + ('...' if len(visual_desc) > 250 else '') if visual_desc else "N/A"
        run_desc.font.size = Pt(12)
        run_desc.font.color.rgb = template.get('text_color', RGBColor(51, 51, 51)) # Use standard text color
        p_desc.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    except Exception as e:
        logging.error(f"Failed to create styled visual placeholder: {e}", exc_info=True)
        # Fallback: Add a simple text box if shape fails
        try:
            # Slightly inset fallback box
            fb_left, fb_top = left + Inches(0.1), top + Inches(0.1)
            fb_width, fb_height = width - Inches(0.2), height - Inches(0.2)
            txBox = slide.shapes.add_textbox(fb_left, fb_top, fb_width, fb_height)
            tf_fb = txBox.text_frame; tf_fb.word_wrap = True; tf_fb.vertical_anchor = MSO_ANCHOR.MIDDLE

            p_fb = tf_fb.add_paragraph()
            # Combine title and description in fallback
            fb_text = f"Visual Suggestion:\n{visual_desc[:250]}{'...' if len(visual_desc) > 250 else ''}"
            p_fb.text = fb_text
            p_fb.font.size = Pt(11)
            p_fb.font.color.rgb = template.get('text_color', RGBColor(51, 51, 51))
            p_fb.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            logging.warning("Added fallback text box for visual suggestion due to shape error.")
        except Exception as fallback_e:
            logging.critical(f"CRITICAL: Failed even to add fallback text box for visual: {fallback_e}")

# --- Flask Routes ---
@app.route('/')
def index():
    """Renders the main upload page."""
    api_key_configured = bool(AZURE_ENDPOINT and AZURE_API_KEY)
    return render_template('index.html', api_key_configured=api_key_configured)

@app.route('/upload', methods=['POST'])
def upload_and_process_file():
    """Handles file upload, processing, and presentation generation."""
    if 'file' not in request.files:
        return jsonify({"error": "No file part in the request."}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No file selected for upload."}), 400

    # --- File Validation ---
    filename = file.filename
    file_ext = os.path.splitext(filename)[1].lower()
    if not allowed_file(filename):
        allowed_str = ", ".join(ALLOWED_EXTENSIONS)
        logging.warning(f"Upload rejected: Invalid file type '{file_ext}' for file '{filename}'. Allowed: {allowed_str}")
        return jsonify({"error": f"Invalid file type '{file_ext}'. Allowed types are: {allowed_str}"}), 400

    # --- Get Form Data ---
    template_name = request.form.get('template', app.config['DEFAULT_TEMPLATE'])
    if template_name not in TEMPLATES:
        logging.warning(f"Invalid template name '{template_name}' received. Falling back to default '{app.config['DEFAULT_TEMPLATE']}'.")
        template_name = app.config['DEFAULT_TEMPLATE']
    target_audience = request.form.get('audience', '').strip()
    desired_tone = request.form.get('tone', '').strip()

    logging.info(f"Processing request for '{filename}' - Template: '{template_name}', Audience: '{target_audience or 'Default'}', Tone: '{desired_tone or 'Default'}'")

    # --- Prepare Filenames and Paths ---
    unique_id = uuid.uuid4().hex[:8]
    # Sanitize original filename for use in output (more robust)
    safe_base_name = re.sub(r'[^\w\-.]+', '_', os.path.splitext(filename)[0]) # Allow letters, numbers, underscore, hyphen, dot
    safe_base_name = re.sub(r'_+', '_', safe_base_name).strip('_') # Consolidate underscores
    safe_base_name = safe_base_name or "document" # Default if sanitization removes everything
    safe_base_name = safe_base_name[:60] # Limit length

    upload_filename = f"{unique_id}_{secure_filename(filename)}" # Secure version for storage
    upload_path = os.path.join(app.config['UPLOAD_FOLDER'], upload_filename)
    pptx_filename = f"{unique_id}_{safe_base_name}_presentation.pptx" # Internal name
    pptx_path = os.path.join(app.config['GENERATED_FOLDER'], pptx_filename)
    # User-friendly download name (without unique ID)
    output_download_name = f"{safe_base_name}_presentation.pptx"

    files_to_cleanup = [] # Keep track of files to delete after request

    @after_this_request
    def cleanup(response):
        """Cleans up temporary files after the request is finished."""
        # Add the generated PPTX path for cleanup *if* it exists
        # This prevents errors if generation failed before saving
        if os.path.exists(pptx_path):
            files_to_cleanup.append(pptx_path)

        for f_path in files_to_cleanup:
            try:
                if os.path.exists(f_path): # Double-check existence before removing
                    os.remove(f_path)
                    logging.info(f"Cleaned up temporary file: {f_path}")
            except OSError as e:
                logging.error(f"Error removing temporary file {f_path}: {e}", exc_info=True)
            except Exception as e:
                 logging.error(f"Unexpected error during cleanup of {f_path}: {e}", exc_info=True)
        return response

    try:
        # 1. Save Uploaded File
        file.save(upload_path)
        files_to_cleanup.append(upload_path) # Ensure upload is always added for cleanup
        logging.info(f"Uploaded file saved temporarily to: {upload_path}")

        # 2. Extract Text
        logging.info(f"Extracting text from '{filename}'...")
        if file_ext == '.docx':
            extracted_text = extract_text_from_docx(upload_path)
        elif file_ext == '.pdf':
            extracted_text = extract_text_from_pdf(upload_path)
        else:
             # Should be caught by initial validation, but acts as a safeguard
             raise ValueError(f"Unsupported file type '{file_ext}' encountered during processing.")

        if not extracted_text or len(extracted_text.strip()) < 50: # Slightly higher threshold for content check
            logging.warning(f"Extracted text from '{filename}' seems very short ({len(extracted_text)} chars). May indicate an empty document or extraction issue.")
            raise ValueError("Document appears to be empty or text could not be extracted properly. Please check the document content.")
        logging.info(f"Text extracted successfully from '{filename}' ({len(extracted_text)} characters).")

        # 3. Build LLM Prompt
        prompt = build_llm_prompt(extracted_text, template_name, target_audience, desired_tone)

        # 4. Call LLM Service
        llm_response = call_llm(prompt)

        # 5. Parse LLM Output
        logging.info("Parsing LLM response...")
        slides_data = parse_llm_output(llm_response)

        # 6. Create PowerPoint Presentation
        logging.info(f"Creating presentation '{output_download_name}' using template '{template_name}'...")
        create_presentation(slides_data, pptx_path, template_name) # This function now raises IOError on save failure

        # 7. Send File to User for Download
        logging.info(f"Presentation created. Preparing '{output_download_name}' for download.")

        # Use BytesIO for efficient memory handling and sending the file
        file_buffer = BytesIO()
        try:
            with open(pptx_path, 'rb') as f:
                file_buffer.write(f.read())
            file_buffer.seek(0) # Rewind buffer to the beginning for sending
        except FileNotFoundError:
            logging.error(f"Generated presentation file not found at expected path: {pptx_path}")
            return jsonify({"error": "Internal server error: Failed to read generated presentation file."}), 500
        except Exception as e:
             logging.error(f"Error reading generated presentation file '{pptx_path}': {e}", exc_info=True)
             return jsonify({"error": "Internal server error: Could not read generated file for download."}), 500

        # Note: Cleanup of pptx_path is handled by the @after_this_request decorator

        return send_file(
            file_buffer,
            as_attachment=True,
            download_name=output_download_name,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except ValueError as e:
        # Handle user-facing errors (validation, empty doc, config issues, parsing, LLM errors)
        logging.warning(f"Processing failed due to ValueError: {e}")
        return jsonify({"error": str(e)}), 400 # Bad Request
    except IOError as e:
        # Handle errors specifically related to file saving (create_presentation)
        logging.error(f"File saving error during presentation creation: {e}", exc_info=True)
        return jsonify({"error": f"Failed to save the presentation file on the server: {e}"}), 500 # Internal Server Error
    except Exception as e:
        # Catch-all for any other unexpected errors during the process
        logging.exception("An unexpected critical error occurred during file processing.") # Log full traceback
        return jsonify({"error": "An internal server error occurred. Please try again later or contact support."}), 500

# --- Main Execution Block ---
if __name__ == '__main__':
    # Check for essential Azure configuration at startup and print warning if missing
    if not AZURE_ENDPOINT or not AZURE_API_KEY:
        print("\n" + "="*75)
        print("!!! WARNING: Azure OpenAI Credentials (Endpoint or API Key) Missing !!!")
        print("!!! The application will run, but PowerPoint generation will fail.  !!!")
        print("!!! Please set AZURE_OPENAI_ENDPOINT and AZURE_OPENAI_API_KEY      !!!")
        print("!!! environment variables or ensure they are in a '.env' file.    !!!")
        print("="*75 + "\n")

    # Determine run mode (Debug vs. Production/Waitress)
    debug_mode = os.environ.get('FLASK_DEBUG', 'False').lower() in ['true', '1', 't']

    host = os.environ.get('HOST', '0.0.0.0') # Listen on all interfaces by default
    port = int(os.environ.get('PORT', 5000)) # Default port 5000

    if not debug_mode:
        try:
            from waitress import serve
            threads = int(os.environ.get('WAITRESS_THREADS', 8)) # Default to 8 threads
            print(f"--- Starting Waitress Production Server ---")
            print(f"--- Listening on: http://{host}:{port}")
            print(f"--- Worker Threads: {threads}")
            print(f"--- Azure Endpoint Configured: {'YES' if AZURE_ENDPOINT else 'NO (!)'}")
            print(f"--- Azure API Key Configured:  {'YES' if AZURE_API_KEY else 'NO (!)'}")
            print("--- Press Ctrl+C to quit ---")
            serve(app, host=host, port=port, threads=threads)
        except ImportError:
            print("\n--- Waitress Package Not Found ---")
            print("--- Falling back to Flask Development Server (NOT FOR PRODUCTION) ---")
            print(f"--- Listening on: http://{host}:{port}")
            print(f"--- Azure Endpoint Configured: {'YES' if AZURE_ENDPOINT else 'NO (!)'}")
            print(f"--- Azure API Key Configured:  {'YES' if AZURE_API_KEY else 'NO (!)'}")
            print("--- Debug Mode: True ---")
            app.run(host=host, port=port, debug=True)
    else:
        print("--- Starting Flask Development Server ---")
        print(f"--- Listening on: http://{host}:{port}")
        print(f"--- Azure Endpoint Configured: {'YES' if AZURE_ENDPOINT else 'NO (!)'}")
        print(f"--- Azure API Key Configured:  {'YES' if AZURE_API_KEY else 'NO (!)'}")
        print("--- Debug Mode: True (Reloads on code changes) ---")
        print("--- Press Ctrl+C to quit ---")
        app.run(host=host, port=port, debug=True)